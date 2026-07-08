"""数据加载模块 - 文档加载、父分块、子分块、向量化、Milvus存储"""
import os
import re
import uuid
import asyncio
import threading
from pathlib import Path
from typing import List, Dict, Tuple
from dataclasses import dataclass, field
import numpy as np
from common.rag.rag_config import RAGConfig
from common.utils.logger import logger


@dataclass
class ParentChunk:
    """父分块 - 按段落分块"""
    id: str
    text: str
    source_file: str
    category: str  # 文件所属分类，如visa/application等
    child_chunks: List["ChildChunk"] = field(default_factory=list)


@dataclass
class ChildChunk:
    """子分块 - 按句子语义分块"""
    id: str
    text: str
    parent_id: str  # 关联的父分块ID
    dense_vector: np.ndarray = None  # bge-m3稠密向量
    sparse_vector: dict = None  # bge-m3稀疏向量


class DocumentLoader:
    """文档加载器 - 支持多种文件格式(pdf, doc, docx, ppt, pptx, html, md, csv, excel, txt)"""
    
    def __init__(self, data_dir: Path = None):
        self.data_dir = data_dir or RAGConfig.DATA_DIR
        self._parsers = {
            ".txt": self._parse_text,
            ".md": self._parse_markdown,
            ".html": self._parse_html,
            ".pdf": self._parse_pdf,
            ".doc": self._parse_doc,
            ".docx": self._parse_docx,
            ".ppt": self._parse_ppt,
            ".pptx": self._parse_pptx,
            ".csv": self._parse_csv,
            ".xls": self._parse_excel,
            ".xlsx": self._parse_excel,
        }
    
    def load_all(self) -> List[Tuple[str, str, str]]:
        """
        加载所有文档
        
        Returns:
            List[(file_path, category, content)] 文件路径、分类、内容
        """
        documents = []
        if not self.data_dir.exists():
            logger.warning(f"数据目录不存在: {self.data_dir}")
            return documents
        
        for category_dir in self.data_dir.iterdir():
            if not category_dir.is_dir():
                continue
            category = category_dir.name
            
            for file_path in category_dir.rglob("*"):
                if not file_path.is_file():
                    continue
                
                ext = file_path.suffix.lower()
                if ext not in RAGConfig.SUPPORTED_FILE_TYPES:
                    continue
                
                try:
                    content = self._load_file(file_path, ext)
                    if content.strip():
                        documents.append((str(file_path), category, content))
                        logger.info(f"加载文档: {file_path.name} (分类: {category}, 格式: {ext})")
                except Exception as e:
                    logger.error(f"加载文档失败 {file_path}: {e}")
        
        logger.info(f"共加载 {len(documents)} 个文档")
        return documents
    
    def _load_file(self, file_path: Path, ext: str) -> str:
        """根据文件扩展名选择解析器"""
        parser = self._parsers.get(ext)
        if parser:
            return parser(file_path)
        return file_path.read_text(encoding="utf-8", errors="ignore")
    
    def _parse_text(self, file_path: Path) -> str:
        """解析纯文本文件"""
        return file_path.read_text(encoding="utf-8", errors="ignore")
    
    def _parse_markdown(self, file_path: Path) -> str:
        """解析Markdown文件"""
        content = file_path.read_text(encoding="utf-8", errors="ignore")
        return content
    
    def _parse_html(self, file_path: Path) -> str:
        """解析HTML文件"""
        try:
            from bs4 import BeautifulSoup
            content = file_path.read_text(encoding="utf-8", errors="ignore")
            soup = BeautifulSoup(content, "html.parser")
            return soup.get_text(separator="\n", strip=True)
        except ImportError:
            logger.warning("未安装beautifulsoup4，使用纯文本解析HTML")
            return self._parse_text(file_path)
    
    def _parse_pdf(self, file_path: Path) -> str:
        """解析PDF文件"""
        try:
            import pdfplumber
            content = ""
            with pdfplumber.open(str(file_path)) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        content += text + "\n\n"
            return content.strip()
        except ImportError:
            logger.warning("未安装pdfplumber，无法解析PDF文件")
            return ""
    
    def _parse_doc(self, file_path: Path) -> str:
        """解析DOC文件"""
        try:
            import textract
            return textract.process(str(file_path)).decode("utf-8", errors="ignore").strip()
        except ImportError:
            logger.warning("未安装textract，无法解析DOC文件")
            return ""
    
    def _parse_docx(self, file_path: Path) -> str:
        """解析DOCX文件"""
        try:
            from docx import Document
            doc = Document(str(file_path))
            content = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            return content.strip()
        except ImportError:
            logger.warning("未安装python-docx，无法解析DOCX文件")
            return ""
    
    def _parse_ppt(self, file_path: Path) -> str:
        """解析PPT文件"""
        logger.warning("PPT文件解析需要特殊工具，建议转换为PPTX格式")
        return ""
    
    def _parse_pptx(self, file_path: Path) -> str:
        """解析PPTX文件"""
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content += shape.text + "\n\n"
            return content.strip()
        except ImportError:
            logger.warning("未安装python-pptx，无法解析PPTX文件")
            return ""
    
    def _parse_csv(self, file_path: Path) -> str:
        """解析CSV文件"""
        try:
            import pandas as pd
            df = pd.read_csv(str(file_path), encoding="utf-8", errors="ignore")
            content = df.to_string(index=False)
            return content
        except ImportError:
            logger.warning("未安装pandas，使用纯文本解析CSV")
            return self._parse_text(file_path)
    
    def _parse_excel(self, file_path: Path) -> str:
        """解析Excel文件"""
        try:
            import pandas as pd
            df = pd.read_excel(str(file_path))
            content = df.to_string(index=False)
            return content
        except ImportError:
            logger.warning("未安装pandas或openpyxl，无法解析Excel文件")
            return ""


class ParentChunker:
    """父分块器 - 按段落分块，保持语义完整性，支持语言检测"""
    
    def __init__(self, separator: str = None):
        self.separator = separator or RAGConfig.PARENT_CHUNK_SEPARATOR
    
    def detect_language(self, text: str) -> str:
        """
        检测文本语言
        
        Args:
            text: 文本内容
            
        Returns:
            语言代码: zh(中文), en(英文), ja(日文), default(默认)
        """
        if not text or len(text) < 10:
            return "default"
        
        chinese_chars = re.findall(r'[\u4e00-\u9fff]', text)
        japanese_chars = re.findall(r'[\u3040-\u30ff\u4e00-\u9fff]', text)
        english_chars = re.findall(r'[a-zA-Z]', text)
        
        total_chars = len(text)
        chinese_ratio = len(chinese_chars) / total_chars if total_chars > 0 else 0
        english_ratio = len(english_chars) / total_chars if total_chars > 0 else 0
        
        if chinese_ratio > 0.3:
            return "zh"
        elif english_ratio > 0.5:
            return "en"
        elif len(japanese_chars) > 0.3 * total_chars:
            return "ja"
        else:
            return "default"
    
    def get_strategy(self, language: str) -> dict:
        """获取指定语言的分块策略"""
        return RAGConfig.CHUNK_STRATEGIES.get(language, RAGConfig.CHUNK_STRATEGIES["default"])
    
    def chunk(self, content: str, source_file: str, category: str) -> List[ParentChunk]:
        """
        将文档按段落分割为父分块
        
        Args:
            content: 文档内容
            source_file: 源文件路径
            category: 分类标签
            
        Returns:
            父分块列表
        """
        language = self.detect_language(content)
        strategy = self.get_strategy(language)
        separator = strategy["paragraph_separator"]
        
        paragraphs = re.split(r'\n\s*\n', content)
        parent_chunks = []
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            chunk_id = str(uuid.uuid4())[:12]
            parent_chunk = ParentChunk(
                id=chunk_id,
                text=para,
                source_file=source_file,
                category=category
            )
            parent_chunk.__dict__["language"] = language
            parent_chunks.append(parent_chunk)
        
        logger.info(f"父分块完成: {len(paragraphs)} 段落 -> {len(parent_chunks)} 个有效分块 (语言: {strategy['name']})")
        return parent_chunks


class ChildChunker:
    """子分块器 - 按句子语义分块，支持语言区分策略"""
    
    def __init__(
        self,
        max_chars: int = None,
        min_chars: int = None
    ):
        self.max_chars = max_chars or RAGConfig.CHILD_CHUNK_MAX_CHARS
        self.min_chars = min_chars or RAGConfig.CHILD_CHUNK_MIN_CHARS
    
    def _get_sentence_pattern(self, separators: list) -> re.Pattern:
        """根据分隔符生成句子分割正则表达式"""
        escaped = [re.escape(s) for s in separators]
        pattern_str = '|'.join([f'(?<={s})' for s in escaped])
        return re.compile(pattern_str)
    
    def chunk(self, parent: ParentChunk) -> List[ChildChunk]:
        """
        将父分块按句子分割为子分块
        
        Args:
            parent: 父分块
            
        Returns:
            子分块列表
        """
        language = getattr(parent, "language", "default")
        strategy = RAGConfig.CHUNK_STRATEGIES.get(language, RAGConfig.CHUNK_STRATEGIES["default"])
        
        max_chars = strategy.get("child_chunk_max_chars", self.max_chars)
        min_chars = strategy.get("child_chunk_min_chars", self.min_chars)
        chunk_overlap = strategy.get("chunk_overlap", 20)
        separators = strategy.get("sentence_separators", ["。", "！", "？", ";", "\n", ".", "!", "?"])
        
        sentence_pattern = self._get_sentence_pattern(separators)
        sentences = sentence_pattern.split(parent.text)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        child_chunks = []
        current_text = ""
        
        for sentence in sentences:
            if not sentence:
                continue
            
            if len(current_text) + len(sentence) > max_chars and current_text:
                if len(current_text) >= min_chars:
                    child_id = str(uuid.uuid4())[:12]
                    child_chunks.append(ChildChunk(
                        id=child_id,
                        text=current_text,
                        parent_id=parent.id
                    ))
                
                if chunk_overlap > 0 and len(current_text) > chunk_overlap:
                    current_text = current_text[-chunk_overlap:] + " " + sentence
                else:
                    current_text = sentence
            else:
                if current_text:
                    current_text += " " + sentence
                else:
                    current_text = sentence
        
        if current_text and len(current_text) >= min_chars:
            child_id = str(uuid.uuid4())[:12]
            child_chunks.append(ChildChunk(
                id=child_id,
                text=current_text,
                parent_id=parent.id
            ))
        
        if len(child_chunks) == 0:
            child_id = str(uuid.uuid4())[:12]
            child_chunks.append(ChildChunk(
                id=child_id,
                text=parent.text,
                parent_id=parent.id
            ))
        
        parent.child_chunks = child_chunks
        logger.debug(f"父分块 {parent.id[:8]} -> {len(child_chunks)} 个子分块 (语言: {strategy['name']})")
        return child_chunks


class EmbeddingModel:
    """bge-m3嵌入模型 - 延迟加载，生成稠密向量和稀疏向量"""
    
    def __init__(
        self,
        model_name: str = None,
        device: str = None
    ):
        self.model_name = model_name or RAGConfig.EMBEDDING_MODEL_NAME
        self.device = device or RAGConfig.EMBEDDING_DEVICE
        self.model = None
        self.tokenizer = None
        self._initialized = False
        self._load_failed = False
        self._init_lock = threading.Lock()
        self._inference_lock = threading.Lock()

    def _ensure_loaded(self):
        """延迟加载嵌入模型，使用双重检查锁定保证只加载一次"""
        if self._initialized or self._load_failed:
            return

        with self._init_lock:
            if self._initialized or self._load_failed:
                return
            try:
                self._load_model()
                self._initialized = True
            except Exception:
                self._load_failed = True
                raise
    
    def _load_model(self):
        """加载bge-m3模型"""
        try:
            import ssl
            
            def _patched_ssl_context(purpose=ssl.Purpose.SERVER_AUTH):
                ctx = ssl.SSLContext(ssl.PROTOCOL_TLS_CLIENT)
                ctx.check_hostname = False
                ctx.verify_mode = ssl.CERT_NONE
                return ctx
            
            ssl.create_default_context = _patched_ssl_context
            ssl._create_default_https_context = _patched_ssl_context
            
            import os
            os.environ["HF_HUB_OFFLINE"] = "1"
            os.environ["TRANSFORMERS_OFFLINE"] = "1"
            os.environ["HF_HUB_DISABLE_TELEMETRY"] = "1"
            os.environ["HF_HUB_DISABLE_SYMLINKS_WARNING"] = "1"
            os.environ["HF_HUB_DISABLE_VERSION_CHECK"] = "1"
            os.environ["CURL_CA_BUNDLE"] = ""
            os.environ["REQUESTS_CA_BUNDLE"] = ""
            os.environ["SSL_CERT_FILE"] = ""
            
            import urllib3
            urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
            
            try:
                import huggingface_hub
                huggingface_hub.constants.HF_HUB_OFFLINE = True
                huggingface_hub.constants.HF_HUB_DISABLE_SYMLINKS_WARNING = True
            except (ImportError, AttributeError):
                pass
            
            logger.info("SSL验证已禁用（嵌入模型Windows环境兼容模式），使用离线模式加载本地缓存模型")
            
            from FlagEmbedding import BGEM3FlagModel
            logger.info(f"加载嵌入模型: {self.model_name} (device={self.device})")
            self.model = BGEM3FlagModel(
                self.model_name,
                use_fp16=True,
                device=self.device
            )
            self.tokenizer = self.model.tokenizer
            logger.info("嵌入模型加载成功")
        except Exception as e:
            logger.error(f"嵌入模型加载失败: {e}")
            raise
    
    def encode_texts(self, texts: List[str], batch_size: int = 32) -> Tuple[np.ndarray, List[dict]]:
        """
        批量编码文本

        Args:
            texts: 文本列表
            batch_size: 批处理大小

        Returns:
            (dense_vectors, sparse_vectors_list)
            dense_vectors: numpy数组, shape=(n, dim)
            sparse_vectors: 列表, 每个元素为 {词id: 权重} 的字典
        """
        if not texts:
            return np.array([]), []

        # 若历史加载已失败，直接返回零向量，避免每次请求重复尝试加载
        if self._load_failed:
            logger.debug("EmbeddingModel历史加载失败，返回零向量")
            zero_dense = np.zeros((len(texts), RAGConfig.EMBEDDING_DIM))
            return zero_dense, [{0: 1e-10} for _ in texts]

        with self._inference_lock:
            # 确保模型已加载（在推理锁内部完成，保证首次并发调用也线程安全）
            self._ensure_loaded()

            # 分批处理
            all_dense = []
            all_sparse = []

            for i in range(0, len(texts), batch_size):
                batch = texts[i:i + batch_size]
                try:
                    result = self.model.encode(
                        batch,
                        batch_size=batch_size,
                        max_length=8192,
                        return_dense=True,
                        return_sparse=True,
                    )

                    dense = result['dense_vecs']
                    all_dense.append(dense)

                    # FlagEmbedding返回lexical_weights作为稀疏向量
                    lexical_weights = result.get('lexical_weights', result.get('sparse_vecs', []))

                    for lex_weight in lexical_weights:
                        sparse_dict = {}
                        if isinstance(lex_weight, dict):
                            for key, val in lex_weight.items():
                                try:
                                    idx = int(key)
                                    weight = float(val)
                                    if weight != 0:
                                        sparse_dict[idx] = weight
                                except (ValueError, TypeError):
                                    continue
                        elif hasattr(lex_weight, 'tocoo'):
                            coo = lex_weight.tocoo()
                            for idx, val in zip(coo.col, coo.data):
                                if val != 0:
                                    sparse_dict[int(idx)] = float(val)
                        elif hasattr(lex_weight, 'nonzero'):
                            for idx in np.nonzero(lex_weight)[0]:
                                val = float(lex_weight[idx])
                                if val != 0:
                                    sparse_dict[int(idx)] = val
                        if not sparse_dict:
                            sparse_dict = {0: 1e-10}
                        all_sparse.append(sparse_dict)

                    logger.debug(f"批处理编码: {len(batch)} 文本, dense shape={dense.shape}")

                except Exception as e:
                    logger.error(f"编码失败: {e}")
                    # 填充零向量
                    zero_dense = np.zeros((len(batch), RAGConfig.EMBEDDING_DIM))
                    all_dense.append(zero_dense)
                    all_sparse.extend([{0: 1e-10} for _ in batch])

            dense_vectors = np.concatenate(all_dense, axis=0)
            return dense_vectors, all_sparse

    # ====== 阶段5异步化改造：异步版encode_texts（保留同步encode_texts不变） ======

    async def async_encode_texts(self, texts: list, batch_size: int = 12) -> tuple:
        """异步版批量编码文本（通过asyncio.to_thread包装PyTorch推理）

        Args:
            texts: 文本列表
            batch_size: 批处理大小

        Returns:
            (dense_vectors, sparse_vectors_list)
            与同步 encode_texts 返回值完全一致
        """
        return await asyncio.to_thread(self.encode_texts, texts, batch_size)

    def warmup(self, queries: List[str] = None):
        """
        模型预热，首次加载并编码若干示例查询

        Args:
            queries: 预热查询文本列表，默认使用内置示例
        """
        if queries is None:
            queries = getattr(RAGConfig, "WARMUP_QUERIES", ["留学申请流程", "签证办理材料"])

        self.encode_texts(queries)
        if self._load_failed:
            logger.warning("EmbeddingModel warmup 失败，模型加载不可用，后续将返回零向量")
        else:
            logger.info(f"EmbeddingModel warmup 完成，编码 {len(queries)} 条查询")


class MilvusManager:
    """Milvus向量数据库管理器"""
    
    def __init__(
        self,
        host: str = None,
        port: int = None,
        database_name: str = None,
        collection_name: str = None
    ):
        self.host = host or RAGConfig.MILVUS_HOST
        self.port = port or RAGConfig.MILVUS_PORT
        self.database_name = database_name or RAGConfig.MILVUS_DATABASE_NAME
        self.collection_name = collection_name or RAGConfig.MILVUS_COLLECTION_NAME
        self.client = None
        self._connect()
    
    def _connect(self):
        """连接Milvus并创建/切换数据库"""
        try:
            from pymilvus import MilvusClient, connections, utility, FieldSchema, CollectionSchema, DataType, Collection
            
            # 使用MilvusClient创建数据库（如果不存在）
            milvus_client = MilvusClient(uri=f"http://{self.host}:{self.port}")
            
            # 获取现有数据库列表
            existing_dbs = milvus_client.list_databases()
            logger.info(f"Milvus现有数据库: {existing_dbs}")
            
            # 创建数据库（如果不存在）
            if self.database_name not in existing_dbs:
                milvus_client.create_database(self.database_name)
                logger.info(f"Milvus数据库创建成功: {self.database_name}")
            
            # 切换到目标数据库
            milvus_client.using_database(self.database_name)
            logger.info(f"已切换到数据库: {self.database_name}")
            
            # 使用旧API进行集合操作（兼容性）
            # alias="default"是pymilvus内部连接标识符，与实际数据库名无关
            # 实际数据库由db_name参数指定为self.database_name
            connections.connect(
                alias="default",
                host=self.host,
                port=str(self.port),
                db_name=self.database_name
            )
            logger.info(f"Milvus连接成功: {self.host}:{self.port} (数据库: {self.database_name})")
            
            self.milvus_client = milvus_client
            self.connections = connections
            self.Collection = Collection
            self.FieldSchema = FieldSchema
            self.CollectionSchema = CollectionSchema
            self.DataType = DataType
            self.utility = utility
            
        except ImportError:
            logger.error("pymilvus未安装, 请运行: pip install pymilvus")
            raise
        except Exception as e:
            logger.error(f"Milvus连接失败: {e}")
            raise
    
    def clean_old_collections(self, exclude_collections=None):
        """清理旧的集合，保留当前项目的集合"""
        exclude_collections = exclude_collections or [self.collection_name]
        
        try:
            collections = self.utility.list_collections()
            for coll in collections:
                if coll not in exclude_collections:
                    logger.warning(f"发现旧集合 {coll}，正在删除...")
                    self.utility.drop_collection(coll)
                    logger.info(f"已删除旧集合: {coll}")
        except Exception as e:
            logger.error(f"清理旧集合失败: {e}")
    
    def ensure_collection(self):
        """确保集合存在，不存在则创建"""
        if self.utility.has_collection(self.collection_name):
            logger.info(f"集合已存在: {self.collection_name}")
            return
        
        # 定义字段
        fields = [
            self.FieldSchema(name="id", dtype=self.DataType.VARCHAR, max_length=64, is_primary=True),
            self.FieldSchema(name="text", dtype=self.DataType.VARCHAR, max_length=65535),  # 子块内容
            self.FieldSchema(name="parent_id", dtype=self.DataType.VARCHAR, max_length=64),
            self.FieldSchema(name="parent_text", dtype=self.DataType.VARCHAR, max_length=65535),  # 父块内容
            self.FieldSchema(name="category", dtype=self.DataType.VARCHAR, max_length=64),
            self.FieldSchema(name="created_at", dtype=self.DataType.INT64),  # 创建时间戳
            # 稠密向量 - bge-m3
            self.FieldSchema(name="dense_vector", dtype=self.DataType.FLOAT_VECTOR, dim=RAGConfig.EMBEDDING_DIM),
            # 稀疏向量 - 使用SPARSE_FLOAT_VECTOR
            self.FieldSchema(name="sparse_vector", dtype=self.DataType.SPARSE_FLOAT_VECTOR),
        ]
        
        schema = self.CollectionSchema(fields, description="RAG检索集合")
        collection = self.Collection(self.collection_name, schema)
        
        # 创建索引
        index_params_dense = {
            "index_type": "IVF_FLAT",
            "metric_type": "IP",  # 内积
            "params": {"nlist": 128}
        }
        collection.create_index("dense_vector", index_params_dense)
        
        index_params_sparse = {
            "index_type": "SPARSE_INVERTED_INDEX",
            "metric_type": "IP",
            "params": {"drop_ratio_build": 0.2}
        }
        collection.create_index("sparse_vector", index_params_sparse)
        
        collection.load()
        logger.info(f"集合创建成功: {self.collection_name}")
    
    def insert_chunks(self, chunks: List[ChildChunk]) -> bool:
        """
        插入子分块到Milvus
        
        Args:
            chunks: 子分块列表(已包含向量)
            
        Returns:
            是否成功
        """
        from pymilvus import Collection
        
        self.ensure_collection()
        collection = self.Collection(self.collection_name)
        
        # 构建插入数据
        entities = []
        for chunk in chunks:
            # 将稀疏字典转换为SparseFloatVec格式
            sparse_data = self._dict_to_sparseFloat(chunk.sparse_vector)
            
            entities.append({
                "id": chunk.id,
                "text": chunk.text,
                "parent_id": chunk.parent_id,
                "source_file": "",  # 需要从父分块获取
                "category": "",     # 需要从父分块获取
                "dense_vector": chunk.dense_vector.tolist() if chunk.dense_vector is not None else [],
                "sparse_vector": sparse_data,
            })
        
        if not entities:
            return False
        
        collection.insert(entities)
        collection.flush()
        logger.info(f"插入 {len(entities)} 条记录到Milvus")
        return True
    
    def insert_with_parent_info(self, chunks: List[ChildChunk], parent_map: Dict[str, ParentChunk], created_at: int = None) -> bool:
        """
        插入子分块并附带父分块信息

        Args:
            chunks: 子分块列表
            parent_map: 父分块映射 {parent_id: ParentChunk}
            created_at: 数据创建时间戳（Unix秒），默认使用当前时间
        """
        import time
        from pymilvus import Collection

        self.ensure_collection()
        collection = self.Collection(self.collection_name)

        if created_at is None:
            created_at = int(time.time())

        entities = []
        for chunk in chunks:
            parent = parent_map.get(chunk.parent_id)
            parent_text = parent.text if parent else ""
            category = parent.category if parent else ""

            sparse_data = self._dict_to_sparseFloat(chunk.sparse_vector)

            entities.append({
                "id": chunk.id,
                "text": chunk.text,
                "parent_id": chunk.parent_id,
                "parent_text": parent_text,
                "category": category,
                "created_at": created_at,
                "dense_vector": chunk.dense_vector.tolist() if chunk.dense_vector is not None else [],
                "sparse_vector": sparse_data,
            })

        if not entities:
            return False

        collection.insert(entities)
        collection.flush()
        logger.info(f"插入 {len(entities)} 条记录到Milvus (含父块内容，时间戳={created_at})")
        return True
    
    def _dict_to_sparseFloat(self, sparse_dict: dict) -> dict:
        """将稀疏字典转换为Milvus接受的格式"""
        if not sparse_dict:
            # Milvus不允许空稀疏向量
            return {0: 1e-10}
        return sparse_dict
    
    def search_dense(self, query_vector: np.ndarray, top_k: int = 10) -> List[dict]:
        """稠密向量检索"""
        from pymilvus import Collection
        
        collection = self.Collection(self.collection_name)
        collection.load()
        
        search_params = {
            "metric_type": "IP",
            "params": {"nprobe": 16}
        }
        
        results = collection.search(
            data=[query_vector.tolist()],
            anns_field="dense_vector",
            param=search_params,
            limit=top_k,
            output_fields=["id", "text", "parent_id", "parent_text", "category", "created_at"]
        )
        
        return self._parse_results(results)
    
    def search_sparse(self, query_sparse: dict, top_k: int = 10) -> List[dict]:
        """稀疏向量检索"""
        from pymilvus import Collection
        
        collection = self.Collection(self.collection_name)
        collection.load()
        
        sparse_vec = self._dict_to_sparseFloat(query_sparse)
        
        search_params = {
            "metric_type": "IP",
            "params": {"drop_ratio_search": 0.2}
        }
        
        results = collection.search(
            data=[sparse_vec],
            anns_field="sparse_vector",
            param=search_params,
            limit=top_k,
            output_fields=["id", "text", "parent_id", "parent_text", "category", "created_at"]
        )
        
        return self._parse_results(results)
    
    def search_hybrid(
        self,
        query_dense: np.ndarray,
        query_sparse: dict,
        dense_weight: float = None,
        sparse_weight: float = None,
        top_k: int = 10
    ) -> List[dict]:
        """
        混合检索 - 稠密+稀疏加权
        
        Args:
            query_dense: 稠密查询向量
            query_sparse: 稀疏查询向量
            dense_weight: 稠密权重
            sparse_weight: 稀疏权重
            top_k: 返回top-k
            
        Returns:
            检索结果列表 [{id, text, parent_id, score, ...}]
        """
        from pymilvus import Collection
        
        collection = self.Collection(self.collection_name)
        collection.load()
        
        dw = dense_weight or RAGConfig.DENSE_VECTOR_WEIGHT
        sw = sparse_weight or RAGConfig.SPARSE_VECTOR_WEIGHT
        
        # 稠密检索
        dense_results = self.search_dense(query_dense, top_k=top_k)
        
        # 稀疏检索
        sparse_results = self.search_sparse(query_sparse, top_k=top_k)
        
        # 合并结果，加权融合
        score_map = {}
        
        for r in dense_results:
            score_map[r["id"]] = {
                "data": r,
                "dense_score": r.get("score", 0) * dw,
                "sparse_score": 0
            }
        
        for r in sparse_results:
            if r["id"] in score_map:
                score_map[r["id"]]["sparse_score"] = r.get("score", 0) * sw
            else:
                score_map[r["id"]] = {
                    "data": r,
                    "dense_score": 0,
                    "sparse_score": r.get("score", 0) * sw
                }
        
        # 计算加权总分
        for id_, info in score_map.items():
            info["data"]["score"] = info["dense_score"] + info["sparse_score"]
            info["data"]["dense_score"] = info["dense_score"]
            info["data"]["sparse_score"] = info["sparse_score"]
        
        # 按总分排序取top-k
        sorted_results = sorted(score_map.values(), key=lambda x: x["data"]["score"], reverse=True)
        return [item["data"] for item in sorted_results[:top_k]]
    
    def get_parent_chunk_by_ids(self, chunk_ids: List[str]) -> Dict[str, dict]:
        """
        通过子分块ID获取关联的父分块信息
        
        Args:
            chunk_ids: 子分块ID列表
            
        Returns:
            {chunk_id: {parent_id, parent_text, ...}}
        """
        from pymilvus import Collection
        
        collection = self.Collection(self.collection_name)
        collection.load()
        
        result_map = {}
        for chunk_id in chunk_ids:
            results = collection.query(
                expr=f"id == '{chunk_id}'",
                output_fields=["id", "parent_id", "text"]
            )
            if results:
                result_map[chunk_id] = results[0]
        
        return result_map
    
    def get_parent_text(self, parent_id: str) -> str:
        """通过parent_id获取父分块文本"""
        from pymilvus import Collection
        
        collection = self.Collection(self.collection_name)
        collection.load()
        
        # 查询该parent_id下的任意一条记录获取父分块文本
        results = collection.query(
            expr=f"parent_id == '{parent_id}'",
            output_fields=["text"],
            limit=1
        )
        
        if results:
            return results[0]["text"]
        return ""
    
    def get_all_child_ids_for_parent(self, parent_id: str) -> List[str]:
        """获取指定父分块的所有子分块ID"""
        from pymilvus import Collection
        
        collection = self.Collection(self.collection_name)
        collection.load()
        
        results = collection.query(
            expr=f"parent_id == '{parent_id}'",
            output_fields=["id"]
        )
        
        return [r["id"] for r in results]
    
    def _parse_results(self, milvus_results) -> List[dict]:
        """解析Milvus检索结果"""
        results = []
        for hits in milvus_results:
            for hit in hits:
                results.append({
                    "id": hit.entity.get("id"),
                    "text": hit.entity.get("text"),
                    "parent_id": hit.entity.get("parent_id"),
                    "parent_text": hit.entity.get("parent_text"),
                    "category": hit.entity.get("category"),
                    "created_at": hit.entity.get("created_at"),
                    "score": float(hit.distance),
                })
        return results
    
    def drop_collection(self):
        """删除集合(用于重建)"""
        if self.utility.has_collection(self.collection_name):
            self.utility.drop_collection(self.collection_name)
            logger.info(f"已删除集合: {self.collection_name}")
    
    def get_count(self) -> int:
        """获取集合中的记录数"""
        from pymilvus import Collection
        
        if not self.utility.has_collection(self.collection_name):
            return 0
        
        collection = self.Collection(self.collection_name)
        return collection.num_entities

    def delete_by_expression(self, expr: str) -> bool:
        """按表达式删除Milvus中的记录"""
        from pymilvus import Collection
        try:
            collection = self.Collection(self.collection_name)
            collection.load()
            collection.delete(expr)
            logger.info(f"Milvus删除: expr={expr}")
            return True
        except Exception as e:
            logger.error(f"Milvus删除失败: {e}")
            return False

    # ====== 阶段5异步化改造：异步版搜索方法（保留同步方法不变） ======

    async def async_search_dense(self, query_dense: np.ndarray, top_k: int = 10) -> list:
        """异步版稠密向量检索（通过asyncio.to_thread包装）

        Args:
            query_dense: 稠密查询向量
            top_k: 返回top-k

        Returns:
            检索结果列表
        """
        return await asyncio.to_thread(self.search_dense, query_dense, top_k)

    async def async_search_sparse(self, query_sparse: dict, top_k: int = 10) -> list:
        """异步版稀疏向量检索（通过asyncio.to_thread包装）

        Args:
            query_sparse: 稀疏查询向量
            top_k: 返回top-k

        Returns:
            检索结果列表
        """
        return await asyncio.to_thread(self.search_sparse, query_sparse, top_k)

    async def async_search_hybrid(
        self,
        query_dense: np.ndarray,
        query_sparse: dict,
        dense_weight: float = None,
        sparse_weight: float = None,
        top_k: int = 10
    ) -> list:
        """异步版混合检索（通过asyncio.to_thread包装）

        Args:
            query_dense: 稠密查询向量
            query_sparse: 稀疏查询向量
            dense_weight: 稠密权重
            sparse_weight: 稀疏权重
            top_k: 返回top-k

        Returns:
            检索结果列表
        """
        return await asyncio.to_thread(
            self.search_hybrid, query_dense, query_sparse, dense_weight, sparse_weight, top_k
        )

    async def async_get_parent_chunk_by_ids(self, chunk_ids: list) -> dict:
        """异步版通过子分块ID获取父分块信息（通过asyncio.to_thread包装）

        Args:
            chunk_ids: 子分块ID列表

        Returns:
            {chunk_id: {parent_id, parent_text, ...}}
        """
        return await asyncio.to_thread(self.get_parent_chunk_by_ids, chunk_ids)
